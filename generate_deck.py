from pptx import Presentation
from pptx.util import Pt

def generate_capstone_deck():
    
    prs = Presentation()
    
    
    title_slide_layout = prs.slide_layouts[0]
    content_slide_layout = prs.slide_layouts[1]

    
    slide1 = prs.slides.add_slide(title_slide_layout)
    slide1.shapes.title.text = "Sentinel-Vault: Zero-Trust JIT IAM"
    slide1.placeholders[1].text = "Architecting Mutual Authentication to Defeat Social Engineering\n\n[ Your Name ]"

    
    slide2 = prs.slides.add_slide(content_slide_layout)
    slide2.shapes.title.text = "The Threat Landscape: Telecom Vulnerability"
    tf2 = slide2.placeholders[1].text_frame
    tf2.text = "The weakest link in cloud security is human psychology paired with legacy telecom."
    
    p = tf2.add_paragraph()
    p.text = "VoIP SIP headers (Caller ID) lack cryptographic verification."
    p.level = 1
    
    p = tf2.add_paragraph()
    p.text = "Attackers exploit humanity's 'habit of trust' to spoof internal IT Helpdesks."
    p.level = 1

    p = tf2.add_paragraph()
    p.text = "Case Study: The 2023 MGM Casino Hack ($100M loss via voice phishing)."
    p.level = 1

    
    slide3 = prs.slides.add_slide(content_slide_layout)
    slide3.shapes.title.text = "The Solution: Sentinel-Vault JIT IAM"
    tf3 = slide3.placeholders[1].text_frame
    tf3.text = "Bypassing the telecom network via Application-Layer Cryptography."
    
    p = tf3.add_paragraph()
    p.text = "Enforces the 'Four-Eyes Principle' for high-risk cloud access."
    p.level = 1
    
    p = tf3.add_paragraph()
    p.text = "[ INSERT 90-SECOND VIDEO DEMO HERE ]"
    p.font.bold = True
    p.level = 1

    
    slide4 = prs.slides.add_slide(content_slide_layout)
    slide4.shapes.title.text = "Enterprise Audit Architecture"
    tf4 = slide4.placeholders[1].text_frame
    tf4.text = "Security requires observability. All events stream to the SOC in real-time."
    
    p = tf4.add_paragraph()
    p.text = "FastAPI Backend: Emits high-fidelity JSON webhooks."
    p.level = 1
    
    p = tf4.add_paragraph()
    p.text = "Apache NiFi: Orchestrates the event routing."
    p.level = 1

    p = tf4.add_paragraph()
    p.text = "Cribl Stream: Masks sensitive PII data-in-transit."
    p.level = 1
    
    p = tf4.add_paragraph()
    p.text = "Snowflake: Maintains the immutable audit log."
    p.level = 1

    
    slide5 = prs.slides.add_slide(content_slide_layout)
    slide5.shapes.title.text = "Business Impact & Future Vision"
    tf5 = slide5.placeholders[1].text_frame
    tf5.text = "Sentinel-Vault successfully enforces 'Least Privilege' and protects data-in-use."
    
    p = tf5.add_paragraph()
    p.text = "Future Roadmap: Packaging Sentinel-Vault as an IDaaS (Identity-as-a-Service) REST API."
    p.level = 1

    p = tf5.add_paragraph()
    p.text = "Seamless integration for enterprise DevSecOps pipelines."
    p.level = 1

    
    prs.save('Sentinel_Vault_Capstone.pptx')
    print("🎯 Success! 'Sentinel_Vault_Capstone.pptx' has been generated.")

if __name__ == "__main__":
    generate_capstone_deck()